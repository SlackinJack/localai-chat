import docx2txt
import os
import speech_recognition
import random
import re
import urllib
import subprocess
import sys

from pptx import Presentation
from pydub import AudioSegment
from PyPDF2 import PdfReader

from modules.utils import *


def getPDFText(filePath):
    pdfFile = PdfReader(filePath)
    pdfPages = len(pdfFile.pages)
    pdfText = ""
    i = 0
    while i < pdfPages:
        pdfText += pdfFile.pages[i].extract_text()
        i += 1
    return pdfText


def getDOCXText(filePath):
    return docx2txt.process(filePath)


def getPPTXText(filePath):
    content = ""
    pres = Presentation(filePath)
    for slide in pres.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                content += shape.text + " "
    return content


def getAudioText(filePath):
    v = AudioSegment.from_file(filePath)
    track = v.set_channels(1).set_frame_rate(16000).set_sample_width(2)
    track.export("temp.wav", format="wav")
    filePathToUse = "temp.wav"
    recog = speech_recognition.Recognizer()
    with speech_recognition.AudioFile("temp.wav") as source:
        src = recog.record(source)
    return recog.recognize_google(src)


def getFileText(filePath):
    f = open(filePath, "r")
    content = f.read()
    f.close()
    return content


fileMap = {
    getPDFText: ["pdf"],
    getDOCXText: ["docx"],
    getPPTXText: ["pptx"],
    getAudioText: ["3gp", "flac", "flv", "mov", "mp3", "mp4", "ogg", "wav"],
}


def getFileContents(filePath):
    k = filePath.split(".")
    fileExtension = k[len(k) - 1]
    content = ""
    for entry, value in fileMap.items():
        for ext in value:
            if fileExtension in ext:
                functionCall = entry
                content = functionCall(filePath)
                return cleanupString(content)
    content = getFileText(filePath)
    return cleanupString(content)


def openLocalFile(filePath):
    if sys.platform == "win32":
        os.startfile(filePath)
    else:
        opener = "xdg-open"
        if sys.platform == "darwin":
            opener = "open"
        subprocess.call([opener, filePath])
